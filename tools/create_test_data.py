import os
from fpdf import FPDF
from pptx import Presentation

def create_pdf(path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="This is Page 1 of the verification PDF.", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(200, 10, txt="It contains distinct information about Alpha.", new_x="LMARGIN", new_y="NEXT")
    
    pdf.add_page()
    pdf.cell(200, 10, txt="This is Page 2 of the verification PDF.", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(200, 10, txt="It talks about Beta features.", new_x="LMARGIN", new_y="NEXT")
    
    pdf.add_page()
    pdf.cell(200, 10, txt="This is Page 3.", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(200, 10, txt="Conclusion on Gamma.", new_x="LMARGIN", new_y="NEXT")
    
    pdf.output(path)
    print(f"Created {path}")

def create_pptx(path):
    prs = Presentation()
    
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Verification Slide 1"
    slide.placeholders[1].text = "Info about Apple."
    
    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Verification Slide 2"
    slide.placeholders[1].text = "Info about Banana."
    
    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Verification Slide 3"
    slide.placeholders[1].text = "Info about Cherry."
    
    prs.save(path)
    print(f"Created {path}")

if __name__ == "__main__":
    os.makedirs("data/source/user", exist_ok=True)
    create_pdf("data/source/user/verify_pages.pdf")
    create_pptx("data/source/user/verify_slides.pptx")
