import os
from pathlib import Path
from pptx import Presentation
from docx import Document as DocxDocument
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from bs4 import BeautifulSoup
from pypdf import PdfReader

def load_pptx(file_path: Path) -> str:
    """Extracts text and links from PPTX with Slide markers."""
    prs = Presentation(file_path)
    text_runs = []
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        slide_text = [f"--- [Slide {slide_num}] ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    p_text = ""
                    for run in paragraph.runs:
                        p_text += run.text
                        if run.hyperlink and run.hyperlink.address:
                            p_text += f" ({run.hyperlink.address})"
                    slide_text.append(p_text)
        text_runs.append("\n".join(slide_text))
    return "\n".join(text_runs)

def load_docx(file_path: Path) -> str:
    """Extracts text and hyperlinks from DOCX."""
    doc = DocxDocument(file_path)
    text = []
    rels = doc.part.rels
    
    for paragraph in doc.paragraphs:
        p_text = paragraph.text
        # Naive: Check for relationship IDs in xml
        link_urls = []
        for child in paragraph._element.iter():
            if child.tag.endswith("hyperlink"):
                try:
                    rid = child.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
                    if rid in rels:
                        link_urls.append(rels[rid].target_ref)
                except:
                    pass
        
        if link_urls:
            p_text += " [Links: " + ", ".join(link_urls) + "]"
        text.append(p_text)
        
    return "\n".join(text)

def load_pdf(file_path: Path) -> str:
    """Extracts text and appends URI annotations per page with Page markers."""
    reader = PdfReader(file_path)
    full_text = []
    for i, page in enumerate(reader.pages):
        page_num = i + 1
        page_text = f"--- [Page {page_num}] ---\n"
        content = page.extract_text() or ""
        
        # Extract links from annotations
        links = []
        if "/Annots" in page:
            for annot in page["/Annots"]:
                obj = annot.get_object()
                if "/A" in obj and "/URI" in obj["/A"]:
                    links.append(obj["/A"]["/URI"])
        
        if links:
            content += "\n[Links found on this page: " + ", ".join(links) + "]"
        
        full_text.append(page_text + content)
    return "\n".join(full_text)

def load_html(file_path: Path) -> str:
    """Extracts text with inline links."""
    with open(file_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')
        
        # Rewrite links to be visible in text
        for a in soup.find_all('a', href=True):
            a.replace_with(f"{a.get_text()} ({a['href']})")
            
        return soup.get_text(separator=' ', strip=True)

def ingest_docs(corpus_name: str):
    """
    Parses raw documents from data/source and saves text to data/interim.
    Supports: .pdf, .docx, .html, .pptx
    """
    base_dir = Path(__file__).parent.parent
    source_dir = base_dir / "data" / "source" / corpus_name
    interim_dir = base_dir / "data" / "interim" / corpus_name

    if not source_dir.exists():
        print(f"Error: Source directory {source_dir} not found.")
        return

    interim_dir.mkdir(parents=True, exist_ok=True)

    print(f"Ingesting documents from {source_dir}...")
    
    files = list(source_dir.glob("*.*"))
    if not files:
        print("  No files found.")
        return

    for file in files:
        output_file = interim_dir / f"{file.stem}.txt"
        
        # Always re-process to catch new link logic
        # if output_file.exists(): continue

        print(f"  Processing {file.name}...")
        try:
            text = ""
            if file.suffix == ".pdf":
                text = load_pdf(file)
            elif file.suffix == ".docx":
                text = load_docx(file)
            elif file.suffix == ".pptx":
                print(f"  [PPTX] Extracting slides from {file.name}...")
                text = load_pptx(file)
            elif file.suffix == ".html":
                text = load_html(file)
            else:
                print(f"  Skipping unsupported file: {file.name}")
                continue

            # Save to interim
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(text)
            
        except Exception as e:
            print(f"  Error processing {file.name}: {e}")

    print("Ingestion complete.")

if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python ingest_docs.py <corpus_name>")
        sys.exit(1)
    ingest_docs(sys.argv[1])
